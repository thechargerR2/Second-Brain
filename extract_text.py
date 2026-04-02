"""
Shared text extraction utilities for Second Brain ingestion.
Extracts readable text from PDFs, DOCX, XLSX, PPTX, and CSV files.
"""

import os
import logging

log = logging.getLogger("extract_text")

# Max chars to store per document (prevent DB bloat)
MAX_CONTENT_LENGTH = 50000


def _pdf_via_pdfplumber(filepath):
    """Try pdfplumber first — best for most PDFs with selectable text."""
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text and text.strip():
                    parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        log.debug(f"pdfplumber failed for {filepath}: {e}")
        return ""


def _pdf_via_pypdf2(filepath):
    """PyPDF2 fallback for PDFs pdfplumber can't handle."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        log.debug(f"PyPDF2 failed for {filepath}: {e}")
        return ""


def _pdf_via_ocr(filepath):
    """OCR fallback for image-based / vector-drawn PDFs."""
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(filepath, dpi=200)
        parts = []
        for i, img in enumerate(images):
            text = pytesseract.image_to_string(img)
            if text and text.strip():
                parts.append(text.strip())
        return "\n".join(parts)
    except Exception as e:
        log.warning(f"OCR extraction failed for {filepath}: {e}")
        return ""


# Minimum chars to consider extraction successful
_MIN_TEXT_THRESHOLD = 50


def extract_from_pdf(filepath):
    """Extract text from a PDF file.

    Strategy: pdfplumber → PyPDF2 → OCR (tesseract).
    """
    for method in (_pdf_via_pdfplumber, _pdf_via_pypdf2):
        text = method(filepath)
        if len(text) >= _MIN_TEXT_THRESHOLD:
            return text

    # Text extractors failed or got very little — fall back to OCR
    log.info(f"Text extractors got <{_MIN_TEXT_THRESHOLD} chars for {filepath}, trying OCR")
    text = _pdf_via_ocr(filepath)
    return text


def extract_from_docx(filepath):
    """Extract text from a DOCX file (paragraphs + tables)."""
    try:
        from docx import Document
        doc = Document(filepath)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        log.warning(f"DOCX extraction failed for {filepath}: {e}")
        return ""


def extract_from_xlsx(filepath):
    """Extract text from an XLSX file (all sheets)."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        log.warning(f"XLSX extraction failed for {filepath}: {e}")
        return ""


def extract_from_csv(filepath):
    """Extract text from a CSV file."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        log.warning(f"CSV extraction failed for {filepath}: {e}")
        return ""


def extract_from_txt(filepath):
    """Extract text from a plain text file."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        log.warning(f"TXT extraction failed for {filepath}: {e}")
        return ""


def extract_from_pptx(filepath):
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)
    except ImportError:
        log.warning("python-pptx not installed — cannot extract PPTX text")
        return ""
    except Exception as e:
        log.warning(f"PPTX extraction failed for {filepath}: {e}")
        return ""


# Map extensions to extractors
EXTRACTORS = {
    ".pdf": extract_from_pdf,
    ".docx": extract_from_docx,
    ".xlsx": extract_from_xlsx,
    ".csv": extract_from_csv,
    ".pptx": extract_from_pptx,
    ".txt": extract_from_txt,
}


def extract_text(filepath):
    """
    Extract text content from a document file.
    Returns the extracted text, or empty string if unsupported/failed.
    Truncates to MAX_CONTENT_LENGTH.
    """
    ext = os.path.splitext(filepath)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return ""

    text = extractor(filepath)
    if text and len(text) > MAX_CONTENT_LENGTH:
        text = text[:MAX_CONTENT_LENGTH] + "\n\n[... truncated ...]"
    return text
